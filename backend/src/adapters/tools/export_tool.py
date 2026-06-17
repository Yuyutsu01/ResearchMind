import os
import sys

def generate_latex_template(title: str, sections: dict) -> str:
    """
    Generate LaTeX source code summarizing the task research details.
    """
    # Safety filter: Only serialize string sections (e.g. raw text summaries) to prevent
    # compiler crashes or rendering errors on nested JSON objects/arrays.
    sections = {k: v for k, v in sections.items() if isinstance(v, str)}
    title_clean = title.replace("_", "\\_").replace("&", "\\&").replace("%", "\\%")
    tex = r"""\documentclass{article}
\usepackage[utf8]{inputenc}
\usepackage{geometry}
\geometry{a4paper, margin=1in}
\title{""" + title_clean + r"""}
\author{AI Research Analyst Platform}
\date{\today}
\begin{document}
\maketitle
"""
    for sec_name, sec_content in sections.items():
        if sec_name.lower() in ["metadata", "plan", "rl_actions", "citations_extracted"]:
            continue
        tex += f"\\section{{{sec_name.capitalize().replace('_', ' ')}}}\n"
        content_clean = sec_content.replace("_", "\\_").replace("&", "\\&").replace("%", "\\%")
        tex += f"{content_clean}\n\n"
    tex += "\\end{document}"
    return tex

def generate_docx_document(title: str, sections: dict, output_path: str):
    """
    Generate a Word document (.docx) using python-docx with local text file fallback.
    """
    # Safety filter: Only serialize string sections to avoid crashing python-docx
    sections = {k: v for k, v in sections.items() if isinstance(v, str)}
    try:
        import docx
        doc = docx.Document()
        doc.add_heading(title, 0)
        
        for sec_name, sec_content in sections.items():
            if sec_name.lower() in ["metadata", "plan", "rl_actions", "citations_extracted"]:
                continue
            doc.add_heading(sec_name.capitalize().replace('_', ' '), level=1)
            doc.add_paragraph(sec_content)
        doc.save(output_path)
    except ImportError:
        # Fallback to plain text
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"=== {title} ===\n\n")
            for k, v in sections.items():
                if k.lower() in ["metadata", "plan", "rl_actions", "citations_extracted"]:
                    continue
                f.write(f"\n# {k.upper()}\n{v}\n")

def generate_pptx_presentation(title: str, sections: dict, output_path: str):
    """
    Generate PowerPoint slides (.pptx) using python-pptx with outline file fallback.
    """
    # Safety filter: Only serialize string sections to avoid crashing python-pptx splits
    sections = {k: v for k, v in sections.items() if isinstance(v, str)}
    try:
        from pptx import Presentation
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "AI Research Analyst Platform Summary"
        
        # Bullet slides per section
        for sec_name, sec_content in sections.items():
            if sec_name.lower() in ["metadata", "plan", "rl_actions", "citations_extracted"]:
                continue
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = sec_name.capitalize().replace('_', ' ')
            
            tf = slide.placeholders[1].text_frame
            # Split section text to find bullet candidates
            bullets = sec_content.split("\n")
            for b in bullets[:6]:
                b_str = b.strip().lstrip("-*•").strip()
                if b_str:
                    p = tf.add_paragraph()
                    p.text = b_str
        prs.save(output_path)
    except ImportError:
        # Fallback to plain text slides
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"=== Presentation Outline: {title} ===\n\n")
            for k, v in sections.items():
                if k.lower() in ["metadata", "plan", "rl_actions", "citations_extracted"]:
                    continue
                f.write(f"Slide: {k.upper()}\n{v[:300]}\n\n")
