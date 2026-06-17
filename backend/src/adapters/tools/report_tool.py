import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, KeepTogether
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

def compile_markdown_report(title: str, sections: dict, references: list, output_path: str) -> str:
    """
    Generates a structured research report in Markdown.
    """
    try:
        content = f"# Research Report: {title}\n\n"
        content += f"*Generated on: {os.path.basename(output_path)}*\n\n"
        
        for sec_title, sec_body in sections.items():
            content += f"## {sec_title.capitalize()}\n\n"
            content += f"{sec_body}\n\n"
            
        if references:
            content += "## References\n\n"
            for ref in references:
                content += f"* {ref}\n"
                
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        return output_path
    except Exception as e:
        return f"Error creating Markdown: {str(e)}"

def compile_pdf_report(title: str, sections: dict, references: list, output_path: str) -> str:
    """
    Generates a clean, publication-grade PDF report using ReportLab.
    """
    try:
        doc = SimpleDocTemplate(output_path, pagesize=letter,
                                rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
        story = []
        styles = getSampleStyleSheet()
        
        # Define clean, modern color scheme (Deep Blue, Charcoal text)
        primary_color = colors.HexColor("#1A365D")
        text_color = colors.HexColor("#2D3748")
        
        title_style = ParagraphStyle(
            'TitleStyle',
            parent=styles['Title'],
            fontName='Helvetica-Bold',
            fontSize=24,
            leading=28,
            textColor=primary_color,
            alignment=0,
            spaceAfter=20
        )
        
        h1_style = ParagraphStyle(
            'H1Style',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=16,
            leading=20,
            textColor=primary_color,
            spaceBefore=15,
            spaceAfter=8,
            keepWithNext=True
        )
        
        body_style = ParagraphStyle(
            'BodyStyle',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
            leading=14,
            textColor=text_color,
            spaceAfter=10
        )
        
        ref_style = ParagraphStyle(
            'RefStyle',
            parent=styles['Normal'],
            fontName='Helvetica-Oblique',
            fontSize=9,
            leading=12,
            textColor=colors.HexColor("#718096"),
            spaceAfter=6
        )
        
        # Add Title
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 10))
        
        # Add Sections
        for sec_title, sec_body in sections.items():
            story.append(Paragraph(sec_title.capitalize(), h1_style))
            # Split body by double linebreaks for separate paragraphs
            paragraphs = sec_body.split("\n\n")
            for p in paragraphs:
                if p.strip():
                    story.append(Paragraph(p.strip().replace("\n", " "), body_style))
            story.append(Spacer(1, 10))
            
        # Add References
        if references:
            story.append(Paragraph("References", h1_style))
            for ref in references:
                story.append(Paragraph(ref, ref_style))
                
        doc.build(story)
        return output_path
    except Exception as e:
        return f"Error creating PDF: {str(e)}"

def compile_docx_report(title: str, sections: dict, references: list, output_path: str) -> str:
    """
    Generates a Word Document using python-docx.
    """
    try:
        doc = Document()
        doc.add_heading(title, level=0)
        
        for sec_title, sec_body in sections.items():
            doc.add_heading(sec_title.capitalize(), level=1)
            paragraphs = sec_body.split("\n\n")
            for p in paragraphs:
                if p.strip():
                    doc.add_paragraph(p.strip().replace("\n", " "))
                    
        if references:
            doc.add_heading("References", level=1)
            for ref in references:
                doc.add_paragraph(ref, style='List Bullet')
                
        doc.save(output_path)
        return output_path
    except Exception as e:
        return f"Error creating Word doc: {str(e)}"

def compile_pptx_report(title: str, sections: dict, references: list, output_path: str) -> str:
    """
    Generates a PowerPoint presentation deck summarizing findings.
    """
    try:
        prs = Presentation()
        
        # Title Slide
        slide_layout = prs.slide_layouts[0] # Title slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Autonomous Research Intelligence Report summary"
        
        # Content Slides
        bullet_layout = prs.slide_layouts[1] # Title and Content
        
        for sec_title, sec_body in sections.items():
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = sec_title.capitalize()
            
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            
            # Add summary bullets
            sentences = [s.strip() for s in re.split(r'\. ', sec_body) if s.strip()]
            for i, sent in enumerate(sentences[:6]): # Limit to top 6 sentences for slide readability
                if i == 0:
                    tf.paragraphs[0].text = sent
                else:
                    p = tf.add_paragraph()
                    p.text = sent
                    p.level = 0
                    
        # References Slide
        if references:
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = "References"
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            for i, ref in enumerate(references[:8]): # Limit to top 8 references
                if i == 0:
                    tf.paragraphs[0].text = ref
                else:
                    p = tf.add_paragraph()
                    p.text = ref
                    p.level = 0
                    
        prs.save(output_path)
        return output_path
    except Exception as e:
        return f"Error creating PowerPoint: {str(e)}"
import re
